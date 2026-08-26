from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_thesis_presentation():
    prs = Presentation()

    # --- Slayt 1: Kapak ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Neuro-Adaptive Kinodynamic Mission Planning\nFramework for UCAVs"
    subtitle.text = "Yüksek Lisans Tez Savunması\nHacettepe Üniversitesi - Bilgisayar Mühendisliği\n\n(Adınız Soyadınız)"

    # --- Slayt 2: Problem ve Motivasyon ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Problem Tanımı ve Motivasyon"
    tf = slide.placeholders[1].text_frame
    tf.text = "İnsansız Savaş Uçaklarında (UCAV) Taktiksel Planlama Darboğazları:"
    p = tf.add_paragraph()
    p.text = "Hız Problemi: Geleneksel geometrik (G-LOS) hesaplamalar çok yavaştır (135+ saniye), anlık tepkiyi engeller."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fiziksel Gürbüzlük Problemi: A*, RRT* ve PSO gibi algoritmalar rüzgarlı ortamlarda uçağın dönüş yarıçapını (R_min) ihlal eder ve radara yakalanır."
    p.level = 1

    # --- Slayt 3: Önerilen Çatı Mimari ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Önerilen Çözüm: Nöro-Adaptif Mimari"
    tf = slide.placeholders[1].text_frame
    tf.text = "Sistem 3 temel hiyerarşik yapıdan oluşmaktadır:"
    p = tf.add_paragraph()
    p.text = "1. ALGI (Perception): Gerçek zamanlı risk tahmini yapan DNN-TRE."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. KARAR (Decision): Duruma göre algoritma seçen RL-Advisor."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. EYLEM (Action): F-16 kinodinamiğini çözen T-GnP ve K-GNP."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[BURAYA ÇATI MİMARİ ŞEMASININ EKRAN GÖRÜNTÜSÜNÜ YAPIŞTIRIN]"
    p.level = 0
    p.font.color.rgb = RGBColor(255, 0, 0) # Kırmızı uyarı metni

    # --- Slayt 4: Algı (DNN-TRE) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. ALGI: DNN-TRE Vekil Modeli"
    tf = slide.placeholders[1].text_frame
    tf.text = "Risk Hesaplamasında Derin Öğrenme Yaklaşımı:"
    p = tf.add_paragraph()
    p.text = "Işın izleme (Ray-casting) yerine arazi maskelemesini öğrenen tam bağlı sinir ağı (MLP)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Log1p dönüşümü ve Softplus aktivasyonu ile yumuşatılmış risk tahmini."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sonuç: 135 saniyeden 51 saniyeye düşüş (S2_Dense senaryosunda 2.6x Speedup)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[BURAYA DNN-TRE vs G-LOS ISI HARİTASI (HEATMAP) GÖRSELİNİ YAPIŞTIRIN]"
    p.level = 0
    p.font.color.rgb = RGBColor(255, 0, 0)

    # --- Slayt 5: Karar (RL-Advisor) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. KARAR: RL-Advisor ve Meta-Policy"
    tf = slide.placeholders[1].text_frame
    tf.text = "Durumsal Farkındalık ve Hiyerarşik Karar Ağacı:"
    p = tf.add_paragraph()
    p.text = "S_fused = αC (Capability) + βO (Opportunity) + γP (Pressure)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "S_fused < 0.33 ise: RL-Pilot (Hızlı, açık arazi)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "S_fused >= 0.62 ise: T-GnP (Gürbüz, yoğun tehdit)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[BURAYA INTERAKTİF S_FUSED PANELİNİN EKRAN GÖRÜNTÜSÜNÜ YAPIŞTIRIN]"
    p.level = 0
    p.font.color.rgb = RGBColor(255, 0, 0)

    # --- Slayt 6: Eylem (Kinematik ve T-GnP) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. EYLEM: Kinodinamik Kısıtlar ve T-GnP"
    tf = slide.placeholders[1].text_frame
    tf.text = "F-16 Fiziksel Uçuş Limitlerinin Ağa Entegrasyonu:"
    p = tf.add_paragraph()
    p.text = "R_min = V^2 / (g * tan(phi_max)) formülü ile dönüş yarıçapı kısıtı."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A* gibi zikzak çizen algoritmalar yerine uçağın dönebileceği kavislerin (Curvature) hesaplanması."
    p.level = 1

    # --- Slayt 7: Bulgular (Monte Carlo) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Simülasyon Bulguları (Stokastik Testler)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Rüzgar Altında 30 İterasyonlu Monte Carlo Analizi:"
    p = tf.add_paragraph()
    p.text = "Geleneksel PSO ve RRT* başarı oranı: %53 - %86"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Önerilen T-GnP başarı oranı: 30/30 (%100 Başarı)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "İzleme Hatası (Tracking Error): Sadece ~24 metre. Toplam Risk: 0.00."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Takas (Trade-off): DNN-TRE rotayı uzatarak aşırı temkinlilik (Conservatism) yaratmıştır."
    p.level = 1

    # --- Slayt 8: Sonuç ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "7. Sonuç ve Gelecek Çalışmalar"
    tf = slide.placeholders[1].text_frame
    tf.text = "Nöro-Adaptif Mimarinin Başarısı:"
    p = tf.add_paragraph()
    p.text = "Hesaplama maliyetlerini milisaniyelere indirirken fiziksel uçuş güvenliğini %100 oranında sağlamıştır."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Gelecek Çalışmalar:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "DNN-TRE'nin doğrudan maliyet yerine 'Heuristic-shaping' olarak kullanılması ve dinamik rota güncellemesi (Online Replanning)."
    p.level = 1

    prs.save('Tez_Savunmasi_Sunumu.pptx')
    print("Sunum başarıyla oluşturuldu: Tez_Savunmasi_Sunumu.pptx")

if __name__ == '__main__':
    create_thesis_presentation()